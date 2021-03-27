import spacy
from pptx import Presentation
import nltk, collections
import re
from collections import defaultdict #DefaultDict is useful in this case cause it creates a new list when the key is missing. It can be an issue if you're intending to use this as read_only, cause it will return an empty list instead of None. But since, we're not calling any values, it's fine
from nltk.corpus import stopwords

master_index = defaultdict(list)

stop_words = set(stopwords.words('english')) 

filepath = str(input('Enter your filepath (or filename if in the same folder): ')) #filepath
pres = Presentation(filepath)
slideCount = 0
indexes = {}

for slide in pres.slides:
    slideCount += 1
    words_in_slide = [slideCount]
#     print('slide number', slideCount)
    for shape in slide.shapes:
        if (shape.has_text_frame):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    # print(run.text)
                    # print("    ---- New Line ----    ")
                    words_in_slide.append(run.text)
    indexes[slideCount] = words_in_slide

for i in range(len(indexes)):
    print(i+1, " - - - ", indexes[i+1])

unwanted_terms = {'\t','(',')','=','.','-',':'}
numbers = set([i for i in '0123456789'])
# alphabets = set([i for i in 'abcdefghijklmnopqrstuvwx'])

def remove_all_instances(main_text, old_term, replacement=""):
    s = main_text.replace(old_term, replacement)
    return s

def remove_words_with_n_characters(tokenised_list, n_characters = 1):
    new_list = tokenised_list.copy()
    for word in tokenised_list:
        if (len(word) == n_characters):
#             print(word, " --- ", len(word))
#             print(len(word))
#             print(word)
            new_list.remove(word)
#             print(words_without_stops)
#         else:
#             print("NOT REJECTED: ",word, " --- ", len(word))
    return new_list

def index_slide(slideNumber):
    slide = indexes[slideNumber]
    slideText = " ".join(slide[1:]).strip()

    for term in unwanted_terms:
        slideText = remove_all_instances(slideText,term)
    
    for number in numbers:
        slideText = remove_all_instances(slideText,number)
    
    word_tokens = nltk.word_tokenize(slideText)
    
    words_without_stops = [w for w in word_tokens if w not in stop_words]
    
    words_without_singularsStops = remove_words_with_n_characters(words_without_stops)

    words_without_singularsStops = [x.lower() for x in words_without_singularsStops]
    
    for word in set(words_without_singularsStops):
        master_index[word].append(slideNumber)

for i in range(len(indexes)):
#     print(i+1, " - - - ", indexes[i+1])
    index_slide(i+1)

output_str=""
with open('index.txt','w') as f:    
        for elem in sorted(master_index):
            try:
                f.write(str(elem) + " :: " + str(master_index[elem]) + "\n")
            except UnicodeEncodeError:
                print(f'{elem}, found on pages {master_index[elem]} was not encoded')